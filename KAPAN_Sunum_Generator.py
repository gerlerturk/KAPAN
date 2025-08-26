from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Yeni bir sunum oluştur
prs = Presentation()

# Genel tema ayarları (arka plan rengi: koyu siyah-mavi gradient simüle etmek için manuel ayar)
# Gerçek gradient için PowerPoint'te düzenleyin; burada basit siyah arka plan kullanıyoruz.

def add_slide_with_title_and_content(slide_layout_index, title_text, content_text, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # Beyaz
    
    # İçerik ekle
    if content_text:
        tf = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(4))
        tf.text_frame.text = content_text
        tf.text_frame.paragraphs[0].font.size = Pt(20)
        tf.text_frame.paragraphs[0].font.color.rgb = RGBColor(211, 211, 211)  # Açık gri
    
    # Görsel ekle (eğer yol varsa)
    if image_path:
        slide.shapes.add_picture(image_path, Inches(9), Inches(1), Inches(4), Inches(5))
    
    return slide

# Slayt 1: Poetik Giriş
add_slide_with_title_and_content(0, "Bir gecede her şeyi kaybetti...\nAilesini, arkadaşlarını, hayallerini...\nSahip olduğu tek şey: Canı.\nVe onun taliplisi çoktu.\nKAPAN\nİntikamın Karanlık Satranç Tahtası",
                                 None, image_path="path/to/istanbul_night.jpg")  # Görsel yolunu değiştirin

# Slayt 2: Proje Bilgileri
content2 = "TÜR: Dram - Aksiyon\nSÜRE: 50 DK\nBÖLÜM: 8\nSENARYO: Zeki Gürdal Karaoğlu\nESER: Emin A. Temel"
add_slide_with_title_and_content(1, "K A P A N", content2, image_path="path/to/fallen_chess.jpg")

# Slayt 3: Hakkında
content3 = "Bir gecede ailesini kaybeden küçük Yusuf, intikam için büyür. İstanbul’un suç imparatorluklarına karşı kurduğu kapanlar, onu aldatanların ve aldatılanların dünyasında bir satranç ustasına dönüştürür. İhanet, dostluk ve aşkın iç içe geçtiği bir dram-aksiyon hikayesi."
add_slide_with_title_and_content(1, "Hakkında", content3, image_path="path/to/child_silhouette.jpg")

# Slayt 4: Genel Hikaye (Katliam ve Kurtuluş)
content4 = "Bir gece, İstanbul’un suç imparatoru Cemal Mavzer ve ailesi, ihanetle katledilir. Küçük oğlu Yusuf, saklambaç oynarken dolaba gizlenerek kurtulur. Kanlar içinde kaçarken, Berlin’e dönen gurbetçi Eyüp’ün arabasına düşer. Eyüp, Yusuf’u büyütür, ama geçmiş peşini bırakmaz. Yıllar sonra bir telefon, Yusuf’u intikam için İstanbul’a çağırır."
add_slide_with_title_and_content(1, "Genel Hikaye - Bölüm 1", content4, image_path="path/to/bloody_scene.jpg")

# Slayt 5: Genel Hikaye (Dönüş ve İlk İntikam)
content5 = "İstanbul’a dönen Yusuf, gizemli bir sesten direktifler alır. Babasının katillerinden birine ulaşır ve ilk intikamını alır. Ses, ona Mirza’yla dost olma görevini verir: Bunun yolu hapistir. Yusuf, suç işleyerek içeri girer. Görüş gününde, sesin sahibi ortaya çıkar: Çocukluk arkadaşı Yasemin."
add_slide_with_title_and_content(1, "Genel Hikaye - Bölüm 2", content5, image_path="path/to/street_running.jpg")

# Slayt 6: Genel Hikaye (Dostluk ve Twist’ler)
content6 = "Yusuf, Mirza’nın güvenini kazanır. Büyük bir firarla hapisten kaçarlar. Yusuf, Mirza’nın ailesine sızar, kapanlar kurar. Ancak Yasemin’in sırrı ortaya çıkar: O, katliam gecesinden kurtulan Gönül’dür. İki yaralı çocuk, intikam için birleşir. Peki, Yusuf en büyük kapanın yemi midir?"
add_slide_with_title_and_content(1, "Genel Hikaye - Bölüm 3", content6, image_path="path/to/chess_pyramid.jpg")

# Slayt 7: Bölüm Hikayeleri (Bölüm 1-2)
content7 = "- Bölüm 1: Cemal Mavzer’in katliamı. Yusuf, dolapta saklanarak kurtulur. Eyüp, onu Berlin’e götürür. Yıllar sonra gizemli telefon: “Kapan kurma sırası sende.”\n- Bölüm 2: Yusuf, İstanbul’a döner. İlk intikamını alır. Mirza’yla dost olmak için hapse girer. Yasemin’in sırrı: O, katliamdan kurtulan Gönül."
add_slide_with_title_and_content(1, "Bölüm Hikayeleri 1-2", content7, image_path="path/to/prison_entry.jpg")

# Slayt 8: Bölüm Hikayeleri (Bölüm 3-4)
content8 = "- Bölüm 3: Yusuf ve Mirza, hapiste saldırılardan kurtulur. Büyük firar ülke gündemine oturur. Yusuf, Mirza’nın ailesine sızar.\n- Bölüm 4: Yusuf, aileleri birbirine düşürür. Mirza’nın hayatını kurtarır; kardeşlik bağı güçlenir."
add_slide_with_title_and_content(1, "Bölüm Hikayeleri 3-4", content8, image_path="path/to/prison_escape.jpg")

# Slayt 9: Bölüm Hikayeleri (Bölüm 5-6)
content9 = "- Bölüm 5: Mirza’nın ölümü, Yusuf’u sarsar. Aileler arasında kan davası başlar.\n- Bölüm 6: Yasemin’in polis olduğu ortaya çıkar (gerçek adı Aslıhan). Yusuf, polisin onu kullandığını öğrenir."
add_slide_with_title_and_content(1, "Bölüm Hikayeleri 5-6", content9, image_path="path/to/street_fight.jpg")

# Slayt 10: Bölüm Hikayeleri (Bölüm 7-8)
content10 = "- Bölüm 7: Eyüp’ün ihaneti açığa çıkar: Mavzer’leri o sattı. Yusuf, gerçeği öğrenir.\n- Bölüm 8: Şehir savaş alanına döner. Yusuf, Eyüp’ü öldürür, aileleri yok eder. Aşk ve intikam galip gelir."
add_slide_with_title_and_content(1, "Bölüm Hikayeleri 7-8", content10, image_path="path/to/betrayal_shadow.jpg")

# Slayt 11: Karakterler (Yusuf ve Yasemin)
content11 = "Yusuf Mavzer (Serhat Farozlu): Katliam gecesiyle çelikleşmiş. Güven problemi var, bağlanmaktan korkar. Berlin’de dövüş eğitimi almış, intikam ateşiyle yanıyor.\nYasemin (Aslıhan): Soğukkanlı, güzel, polis kimlikli. Katliamdan kurtulan Gönül’ü oynuyor. Aşk ve görev arasında sıkışır."
add_slide_with_title_and_content(1, "Karakterler 1/2", content11, image_path="path/to/yusuf_portrait.jpg")

# Slayt 12: Karakterler (Mirza ve Eyüp) ve Kapanış
content12 = "Mirza: Yara izli, acıya duyarsız, paranoyak. Ailenin veliahdı, Yusuf’un dostu ama kurbanı.\nEyüp: Görünüşte baba figürü, gerçekte hain. Mavzer’leri satan adam.\nKAPAN: İstanbul’un karanlık yüzünü aydınlatacak bir intikam destanı!"
add_slide_with_title_and_content(1, "Karakterler 2/2 ve Kapanış", content12, image_path="path/to/mirza_portrait.jpg")

# Sunumu kaydet
prs.save('KAPAN_Sunum.pptx')
print("PPTX dosyası oluşturuldu: KAPAN_Sunum.pptx")
